from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Slide content example
slides = [
    ("Quarter-End Presentation", "Team: AI Testing\nCompany: BugRaptors Infotech Pvt Ltd"),
    ("Completed Milestones", "• Python Course Completed\n• AI Testing PoC Executed\n• Stakeholder Demos"),
    ("Challenges", "Technical:\n• Model variability\nNon-Technical:\n• Coordination issues"),
    ("Thank You", "Looking forward to your feedback!")
]

# Slide background and font colors
background_color = RGBColor(178, 34, 34)  # Red
title_color = RGBColor(255, 255, 255)
text_color = RGBColor(255, 235, 235)

# Add slides
for title, content in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background_color

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = title_color

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    frame = textbox.text_frame
    para = frame.add_paragraph()
    para.text = content
    para.font.size = Pt(20)
    para.font.color.rgb = text_color
    para.alignment = PP_ALIGN.LEFT

# Save to file
prs.save("Quarter_End_AI_Testing_Presentation.pptx")
